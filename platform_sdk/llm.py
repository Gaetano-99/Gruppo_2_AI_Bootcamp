# ============================================================================
# LearnAI Platform SDK — Modulo LLM
# Funzioni semplificate per interagire con Claude su AWS Bedrock.
#
# Funzioni disponibili:
#   chiedi(prompt)                          → risposta testuale semplice
#   chiedi_con_contesto(domanda, contesto)  → risposta con contesto (RAG-like)
#   chiedi_veloce(prompt)                   → risposta veloce ed economica
#   genera_json(prompt, schema)             → output JSON strutturato
#   chat(messaggi, system_prompt)           → conversazione multi-turn
#   chat_stream(messaggi, system_prompt)    → conversazione con streaming
#   analizza_documento(testo, istruzioni)   → analisi di testo lungo
#   estrai_testo_da_upload(uploaded_file)   → legge un file caricato (txt, pdf, csv, xls, xlsx)
# ============================================================================

import json
import re
from typing import Generator

from langchain_aws import ChatBedrockConverse
from langchain_core.messages import HumanMessage, SystemMessage, AIMessage

import config


# ---------------------------------------------------------------------------
# Inizializzazione dei modelli (lazy — creati alla prima chiamata)
# ---------------------------------------------------------------------------

_llm_principale = None
_llm_veloce = None


def _get_llm_principale():
    """Restituisce il modello principale (Sonnet). Creato una sola volta."""
    global _llm_principale
    if _llm_principale is None:
        _llm_principale = ChatBedrockConverse(
            model=config.BEDROCK_MODEL_ID,
            region_name=config.AWS_DEFAULT_REGION,
            aws_access_key_id=config.AWS_ACCESS_KEY_ID,
            aws_secret_access_key=config.AWS_SECRET_ACCESS_KEY,
            temperature=config.LLM_TEMPERATURE,
            max_tokens=config.LLM_MAX_TOKENS,
        )
    return _llm_principale


def _get_llm_veloce():
    """Restituisce il modello veloce (Haiku). Creato una sola volta."""
    global _llm_veloce
    if _llm_veloce is None:
        _llm_veloce = ChatBedrockConverse(
            model=config.BEDROCK_MODEL_ID_FAST,
            region_name=config.AWS_DEFAULT_REGION,
            aws_access_key_id=config.AWS_ACCESS_KEY_ID,
            aws_secret_access_key=config.AWS_SECRET_ACCESS_KEY,
            temperature=config.LLM_TEMPERATURE,
            max_tokens=config.LLM_MAX_TOKENS_FAST,
        )
    return _llm_veloce


def get_llm(veloce: bool = False):
    """
    Restituisce l'oggetto LLM di LangChain, utile se volete usarlo
    direttamente con LangGraph o per creare agenti custom.

    Parametri:
        veloce: se True, usa il modello Haiku (più economico)

    Esempio:
        llm = get_llm()
        llm_haiku = get_llm(veloce=True)
    """
    return _get_llm_veloce() if veloce else _get_llm_principale()


# ---------------------------------------------------------------------------
# Funzioni principali
# ---------------------------------------------------------------------------


def chiedi(prompt: str) -> str:
    """
    Fa una domanda semplice all'AI e ottiene una risposta testuale.

    Parametri:
        prompt: la domanda o istruzione da inviare all'AI

    Ritorna:
        La risposta dell'AI come stringa di testo

    Esempio:
        risposta = chiedi("Quali competenze servono per un data analyst?")
        print(risposta)
    """
    llm = _get_llm_principale()
    messaggi = [HumanMessage(content=prompt)]
    risposta = llm.invoke(messaggi)
    return risposta.content


def chiedi_veloce(prompt: str) -> str:
    """
    Come chiedi(), ma usa il modello veloce (Haiku).
    Ideale per task semplici: classificazioni, sì/no, brevi risposte.

    Parametri:
        prompt: la domanda o istruzione

    Ritorna:
        La risposta dell'AI come stringa

    Esempio:
        categoria = chiedi_veloce("Classifica questa skill: 'Python'. Rispondi solo: Technical, Soft o Domain")
    """
    llm = _get_llm_veloce()
    messaggi = [HumanMessage(content=prompt)]
    risposta = llm.invoke(messaggi)
    return risposta.content


def chiedi_con_contesto(domanda: str, contesto: str, istruzioni: str = "") -> str:
    """
    Fa una domanda all'AI fornendo un contesto aggiuntivo.
    Simile a un approccio RAG semplificato.

    Parametri:
        domanda:    la domanda dell'utente
        contesto:   informazioni aggiuntive (dati dal DB, documenti, profilo utente...)
        istruzioni: istruzioni opzionali per l'AI (es: "Rispondi in italiano, max 3 frasi")

    Ritorna:
        La risposta dell'AI come stringa

    Esempio:
        risposta = chiedi_con_contesto(
            domanda="Quali corsi mi consigli?",
            contesto="L'utente è un junior analyst con competenze in Excel e SQL base.",
            istruzioni="Suggerisci massimo 3 corsi dal catalogo. Rispondi in italiano."
        )
    """
    system_text = "Sei un assistente AI della piattaforma LearnAI."
    if istruzioni:
        system_text += f"\n\nIstruzioni: {istruzioni}"

    prompt_completo = f"""Contesto:
{contesto}

Domanda:
{domanda}"""

    llm = _get_llm_principale()
    messaggi = [
        SystemMessage(content=system_text),
        HumanMessage(content=prompt_completo),
    ]
    risposta = llm.invoke(messaggi)
    return risposta.content


def analizza_documento(testo_documento: str, istruzioni: str) -> str:
    """
    Analizza un testo lungo (es: CV, documento, report) seguendo delle istruzioni.

    Parametri:
        testo_documento: il testo completo da analizzare
        istruzioni:      cosa deve fare l'AI con il testo

    Ritorna:
        La risposta dell'AI come stringa

    Esempio:
        risultato = analizza_documento(
            testo_documento="Mario Rossi, laureato in economia...",
            istruzioni="Estrai: nome, cognome, competenze, esperienze. Rispondi in JSON."
        )
    """
    llm = _get_llm_principale()
    messaggi = [
        SystemMessage(content=f"Sei un analista esperto. {istruzioni}"),
        HumanMessage(content=f"Analizza il seguente documento:\n\n{testo_documento}"),
    ]
    risposta = llm.invoke(messaggi)
    return risposta.content


def genera_json(prompt: str, schema: dict, istruzioni: str = "") -> dict:
    """
    Chiede all'AI di generare una risposta in formato JSON strutturato.
    Lo schema indica la struttura attesa dell'output.

    Parametri:
        prompt:      la richiesta da inviare all'AI
        schema:      dizionario che descrive la struttura JSON attesa
        istruzioni:  istruzioni aggiuntive opzionali

    Ritorna:
        Un dizionario Python con la struttura richiesta

    Esempio:
        piano = genera_json(
            prompt="Crea un piano formativo per un junior data analyst",
            schema={
                "corsi": [{"titolo": "str", "durata_ore": "int", "priorita": "str"}],
                "durata_totale_settimane": "int"
            },
            istruzioni="Rispondi in italiano"
        )
        print(piano["corsi"][0]["titolo"])
    """
    schema_str = json.dumps(schema, indent=2, ensure_ascii=False)

    system_text = """Sei un assistente AI della piattaforma LearnAI.
DEVI rispondere ESCLUSIVAMENTE con un oggetto JSON valido.
NON aggiungere testo prima o dopo il JSON.
NON usare ```json o altri wrapper."""

    if istruzioni:
        system_text += f"\n\nIstruzioni aggiuntive: {istruzioni}"

    prompt_completo = f"""{prompt}

Rispondi con un JSON che segue ESATTAMENTE questa struttura:
{schema_str}"""

    llm = _get_llm_principale()
    messaggi = [
        SystemMessage(content=system_text),
        HumanMessage(content=prompt_completo),
    ]
    risposta = llm.invoke(messaggi)

    # Parsing della risposta JSON
    return _parse_json_risposta(risposta.content)


def chat(messaggi: list[dict], system_prompt: str = "") -> str:
    """
    Conversazione multi-turn. Mantiene la storia dei messaggi.

    Parametri:
        messaggi:      lista di dizionari con chiavi "role" e "content"
                       role può essere: "user" o "assistant"
        system_prompt: istruzioni di sistema per l'AI (opzionale)

    Ritorna:
        La risposta dell'AI come stringa

    Esempio:
        storia = [
            {"role": "user", "content": "Ciao! Sono Marco, nuovo in azienda."},
            {"role": "assistant", "content": "Benvenuto Marco! Come posso aiutarti?"},
            {"role": "user", "content": "Vorrei capire quali corsi seguire."},
        ]
        risposta = chat(storia, system_prompt="Sei il chatbot di onboarding di LearnAI.")
        print(risposta)
    """
    # Converti i messaggi nel formato LangChain
    messaggi_lc = []

    if system_prompt:
        messaggi_lc.append(SystemMessage(content=system_prompt))

    for msg in messaggi:
        if msg["role"] == "user":
            messaggi_lc.append(HumanMessage(content=msg["content"]))
        elif msg["role"] == "assistant":
            messaggi_lc.append(AIMessage(content=msg["content"]))

    llm = _get_llm_principale()
    risposta = llm.invoke(messaggi_lc)
    return risposta.content


def chat_stream(messaggi: list[dict], system_prompt: str = "") -> Generator[str, None, None]:
    """
    Come chat(), ma restituisce la risposta un pezzo alla volta (streaming).
    Ideale per mostrare la risposta in tempo reale su Streamlit.

    Parametri:
        messaggi:      lista di dizionari con "role" e "content"
        system_prompt: istruzioni di sistema per l'AI

    Ritorna:
        Un generatore che produce pezzi di testo uno alla volta

    Esempio con Streamlit:
        # Nel file Streamlit
        with st.chat_message("assistant"):
            risposta_completa = st.write_stream(
                chat_stream(storia, system_prompt="Sei un tutor AI.")
            )
    """
    # Converti messaggi nel formato LangChain
    messaggi_lc = []

    if system_prompt:
        messaggi_lc.append(SystemMessage(content=system_prompt))

    for msg in messaggi:
        if msg["role"] == "user":
            messaggi_lc.append(HumanMessage(content=msg["content"]))
        elif msg["role"] == "assistant":
            messaggi_lc.append(AIMessage(content=msg["content"]))

    llm = _get_llm_principale()
    for chunk in llm.stream(messaggi_lc):
        if chunk.content:
            # Il contenuto può arrivare come stringa o come lista di blocchi
            if isinstance(chunk.content, str):
                yield chunk.content
            elif isinstance(chunk.content, list):
                for blocco in chunk.content:
                    if isinstance(blocco, dict) and blocco.get("type") == "text":
                        yield blocco["text"]
                    elif isinstance(blocco, str):
                        yield blocco


# ---------------------------------------------------------------------------
# Utility
# ---------------------------------------------------------------------------


def _ocr_pdf_con_claude(dati: bytes) -> str:
    """Estrae testo da un PDF scansionato usando Claude Vision via AWS Bedrock.

    Usa PyMuPDF per renderizzare ogni pagina come immagine JPEG, poi invia
    ogni immagine a Claude (già configurato nel progetto) per l'estrazione
    del testo tramite visione artificiale.

    Parametri:
        dati: Contenuto grezzo del file PDF in byte.

    Ritorna:
        Testo estratto via OCR, oppure stringa vuota se l'operazione fallisce.
    """
    import fitz  # PyMuPDF
    import base64

    doc = fitz.open(stream=dati, filetype="pdf")
    llm = _get_llm_veloce()

    pagine_testo: list[str] = []
    for i, page in enumerate(doc, start=1):
        # Zoom 2x ≈ 144 DPI — buona qualità OCR con immagini < 1 MB
        pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0))
        img_b64 = base64.b64encode(pix.tobytes("jpeg")).decode("utf-8")

        messaggio = HumanMessage(content=[
            {
                "type": "image_url",
                "image_url": {"url": f"data:image/jpeg;base64,{img_b64}"},
            },
            {
                "type": "text",
                "text": (
                    "Estrai tutto il testo presente in questa immagine. "
                    "Mantieni la struttura originale (titoli, paragrafi, elenchi). "
                    "Rispondi solo con il testo estratto, senza commenti aggiuntivi."
                ),
            },
        ])

        risposta = llm.invoke([messaggio])
        testo_pagina = risposta.content.strip() if risposta.content else ""
        if testo_pagina:
            pagine_testo.append(f"[Pagina {i}]\n{testo_pagina}")

    doc.close()
    return "\n\n".join(pagine_testo).strip()


def estrai_testo_da_upload(uploaded_file) -> str:
    """
    Legge il contenuto di un file caricato tramite st.file_uploader.
    Supporta file .txt, .md, .csv, .pdf, .xls, .xlsx, .docx e .pptx.

    Per i PDF usa pdfplumber (gestisce layout complessi e tabelle); se una pagina
    non produce testo (es. pagina con sole immagini) viene aggiunto un segnaposto.
    Per i DOCX estrae paragrafi e celle di tabella. Per i PPTX estrae il testo
    da tutti gli shape di ogni slide.

    Parametri:
        uploaded_file: l'oggetto restituito da st.file_uploader()

    Ritorna:
        Il contenuto del file come stringa

    Esempio:
        file = st.file_uploader("Carica un file", type=["txt", "pdf", "csv", "xls", "xlsx", "docx", "pptx"])
        if file:
            testo = estrai_testo_da_upload(file)
            st.write(testo)
    """
    if uploaded_file is None:
        return ""

    import io
    nome = uploaded_file.name.lower()
    dati = uploaded_file.read()

    # --- PDF ---
    if nome.endswith(".pdf"):
        # Tentativo 1: pdfplumber (migliore per layout complessi e colonne)
        _pagine_pdfplumber: list[str] | None = None
        try:
            import pdfplumber
            _pagine_tmp: list[str] = []
            with pdfplumber.open(io.BytesIO(dati)) as pdf:
                for i, pagina in enumerate(pdf.pages, start=1):
                    testo_pagina = pagina.extract_text() or ""
                    if testo_pagina.strip():
                        _pagine_tmp.append(testo_pagina)
                    else:
                        _pagine_tmp.append(f"[Pagina {i}: contenuto non testuale — possibile immagine o scansione]")
            _pagine_pdfplumber = _pagine_tmp
        except Exception:
            pass

        if _pagine_pdfplumber is not None:
            testo = "\n\n".join(_pagine_pdfplumber).strip()
            _solo_immagini = bool(_pagine_pdfplumber) and all(
                p.startswith("[Pagina") and "non testuale" in p
                for p in _pagine_pdfplumber
            )
            if _solo_immagini:
                # OCR via Claude Vision — le eccezioni qui vengono propagate
                # così docente.py può mostrarle all'utente
                return _ocr_pdf_con_claude(dati)
            if testo:
                return testo

        # Fallback: PyPDF2
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(io.BytesIO(dati))
            pagine = [p.extract_text() or "" for p in reader.pages]
            return "\n\n".join(pagine).strip()
        except Exception as e:
            return f"[Errore lettura PDF: {e}]"

    # --- DOCX ---
    if nome.endswith(".docx"):
        try:
            from docx import Document
            doc = Document(io.BytesIO(dati))
            parti: list[str] = []

            # Paragrafi del corpo principale
            for para in doc.paragraphs:
                testo_para = para.text.strip()
                if testo_para:
                    parti.append(testo_para)

            # Testo nelle tabelle
            for tabella in doc.tables:
                for riga in tabella.rows:
                    testo_riga = " | ".join(
                        cella.text.strip() for cella in riga.cells if cella.text.strip()
                    )
                    if testo_riga:
                        parti.append(testo_riga)

            return "\n\n".join(parti).strip()
        except Exception as e:
            return f"[Errore lettura DOCX: {e}]"

    # --- PPTX ---
    if nome.endswith(".pptx"):
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(dati))
            slide_testi: list[str] = []

            for idx, slide in enumerate(prs.slides, start=1):
                testi_slide: list[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            testo_para = "".join(run.text for run in para.runs).strip()
                            if testo_para:
                                testi_slide.append(testo_para)
                if testi_slide:
                    slide_testi.append(f"[Slide {idx}]\n" + "\n".join(testi_slide))

            return "\n\n".join(slide_testi).strip()
        except Exception as e:
            return f"[Errore lettura PPTX: {e}]"

    # --- CSV ---
    if nome.endswith(".csv"):
        try:
            import pandas as pd
            df = pd.read_csv(io.BytesIO(dati))
            return df.to_string(index=False)
        except Exception as e:
            return f"[Errore lettura CSV: {e}]"

    # --- Excel (.xls, .xlsx) ---
    if nome.endswith((".xls", ".xlsx")):
        try:
            import pandas as pd
            df = pd.read_excel(io.BytesIO(dati))
            return df.to_string(index=False)
        except Exception as e:
            return f"[Errore lettura Excel: {e}]"

    # --- File di testo (.txt, .md, e altri) ---
    if isinstance(dati, bytes):
        dati = dati.decode("utf-8", errors="replace")

    return dati


def _parse_json_risposta(testo: str) -> dict:
    """
    Prova a estrarre un oggetto JSON da una risposta dell'AI.
    Gestisce i casi in cui l'AI aggiunge testo extra attorno al JSON.
    """
    # Prova il parsing diretto
    try:
        return json.loads(testo)
    except json.JSONDecodeError:
        pass

    # Prova a trovare un blocco ```json ... ```
    match = re.search(r"```(?:json)?\s*\n?(.*?)\n?\s*```", testo, re.DOTALL)
    if match:
        try:
            return json.loads(match.group(1))
        except json.JSONDecodeError:
            pass

    # Prova a trovare il primo { ... } o [ ... ]
    match = re.search(r"(\{.*\}|\[.*\])", testo, re.DOTALL)
    if match:
        try:
            return json.loads(match.group(1))
        except json.JSONDecodeError:
            pass

    # Se tutto fallisce, restituisci un dict con il testo originale
    return {"errore": "Impossibile parsare il JSON", "risposta_grezza": testo}
